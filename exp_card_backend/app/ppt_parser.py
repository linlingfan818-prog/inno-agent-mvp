from __future__ import annotations

import base64
from pathlib import Path
from typing import Any

from pptx import Presentation

from .config import settings
from .llm_client import extract_json_from_text, chat_completion, vision_completion
from .prompts import PARSER_SYSTEM_PROMPT
from .schemas import ParsedCharter


KEYWORD_MAP = {
    "scope": ["project scope", "scope", "项目范围", "项目概述"],
    "okr": ["okrs", "okr", "objectives", "目标"],
    "kr": ["key results", "krs", "关键结果"],
    "milestone": ["milestone", "milestones", "里程碑"],
    "value": ["value vs. cost", "value proposition", "价值主张", "收益"],
    "cost": ["cost", "budget", "预算", "成本"],
    "stakeholder": ["stakeholder", "stakeholders", "干系人"],
}


class ParseWarning(RuntimeError):
    pass


def _normalize_line(line: str) -> str:
    return " ".join(line.replace("•", " ").replace("\u2022", " ").split()).strip()


def extract_text_from_pptx(file_path: Path) -> str:
    prs = Presentation(str(file_path))
    if not prs.slides:
        return ""

    slide = prs.slides[0]
    items: list[tuple[int, int, str]] = []
    for shape in slide.shapes:
        text = ""
        if hasattr(shape, "text") and shape.text:
            text = shape.text
        elif getattr(shape, "has_table", False):
            table_lines: list[str] = []
            for row in shape.table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_lines.append(row_text)
            text = "\n".join(table_lines)

        text = "\n".join(_normalize_line(x) for x in text.splitlines() if _normalize_line(x))
        if text:
            items.append((int(shape.top), int(shape.left), text))

    items.sort(key=lambda x: (x[0], x[1]))
    return "\n".join(text for _, _, text in items)


def _image_file_to_data_url(file_path: Path) -> str:
    ext = file_path.suffix.lower().replace(".", "")
    mime = "image/png" if ext == "png" else f"image/{ext}"
    data = base64.b64encode(file_path.read_bytes()).decode("utf-8")
    return f"data:{mime};base64,{data}"


def heuristic_parse(raw_text: str) -> ParsedCharter:
    lines = [_normalize_line(x) for x in raw_text.splitlines() if _normalize_line(x)]
    parsed: dict[str, Any] = {
        "project_name": lines[0] if lines else "",
        "project_scope": "",
        "stakeholders_now": [],
        "stakeholders_future": [],
        "objectives": [],
        "key_results": [],
        "milestones": [],
        "value_points": [],
        "cost_points": [],
        "raw_text": raw_text,
        "parse_notes": "Heuristic parse only.",
    }

    current_section = None
    for line in lines[1:]:
        lower = line.lower()
        switched = False
        for section, keywords in KEYWORD_MAP.items():
            if any(keyword in lower for keyword in keywords):
                current_section = section
                switched = True
                break
        if switched:
            continue

        if current_section == "scope":
            parsed["project_scope"] += (" " + line).strip()
        elif current_section == "stakeholder":
            if "now" in lower or "当前" in lower:
                parsed["stakeholders_now"].append(line)
            elif "future" in lower or "未来" in lower:
                parsed["stakeholders_future"].append(line)
            else:
                parsed["stakeholders_now"].append(line)
        elif current_section == "okr":
            parsed["objectives"].append(line)
        elif current_section == "kr":
            parsed["key_results"].append(line)
        elif current_section == "milestone":
            parsed["milestones"].append(line)
        elif current_section == "value":
            parsed["value_points"].append(line)
        elif current_section == "cost":
            parsed["cost_points"].append(line)

    return ParsedCharter(**parsed)


def llm_parse_from_text(raw_text: str) -> ParsedCharter:
    user_prompt = f"请把以下项目企划书内容解析成 JSON：\n\n{raw_text}"
    response = chat_completion(
        [
            {"role": "system", "content": PARSER_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ]
    )
    payload = extract_json_from_text(response)
    payload["raw_text"] = raw_text
    return ParsedCharter(**payload)


def llm_parse_from_image(file_path: Path) -> ParsedCharter:
    prompt = f"{PARSER_SYSTEM_PROMPT}\n\n请直接识别图片中的一页项目企划书并输出 JSON。"
    response = vision_completion(prompt, _image_file_to_data_url(file_path))
    payload = extract_json_from_text(response)
    payload.setdefault("raw_text", "")
    return ParsedCharter(**payload)


def parse_uploaded_file(file_path: Path) -> tuple[str, ParsedCharter, list[str]]:
    warnings: list[str] = []
    suffix = file_path.suffix.lower()
    raw_text = ""

    if suffix == ".pptx":
        raw_text = extract_text_from_pptx(file_path)
        if not raw_text:
            raise ParseWarning("PPTX 第一页没有提取到文本。")
        try:
            parsed = llm_parse_from_text(raw_text)
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"LLM 结构化解析失败，已回退到规则解析：{exc}")
            parsed = heuristic_parse(raw_text)
        return raw_text, parsed, warnings

    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        if settings.support_image_input:
            try:
                parsed = llm_parse_from_image(file_path)
                return parsed.raw_text or "", parsed, warnings
            except Exception as exc:  # noqa: BLE001
                warnings.append(f"图片视觉解析失败：{exc}")
        raise ParseWarning(
            "当前配置不支持图片直读，建议上传 .pptx 文件，或在 config 中开启 llm.support_image_input。"
        )

    raise ParseWarning("仅支持 .pptx / .png / .jpg / .jpeg / .webp")
